# -*- coding: utf-8 -*-
"""SIRIAI 발표자료 PPTX 빌더 — HTML 덱(SIRIAI_발표자료.html) 미러링."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
INK   = RGBColor(0x21,0x1A,0x33)
CREAM = RGBColor(0xF3,0xEE,0xE2)
DARK  = RGBColor(0x1A,0x14,0x26)
PAPER = RGBColor(0xFB,0xF9,0xF3)
LINE  = RGBColor(0xD7,0xD1,0xC6)
BLUE  = RGBColor(0x2A,0x6F,0xDB)
ORANGE= RGBColor(0xE6,0x51,0x00)
GREEN = RGBColor(0x2E,0x7D,0x32)
RED   = RGBColor(0xC6,0x28,0x28)
PURPLE= RGBColor(0x6B,0x4D,0xCE)
PINK  = RGBColor(0xC1,0x43,0x6D)
PEACH = RGBColor(0xF6,0xC9,0xA8)
WHITE = RGBColor(0xFF,0xFF,0xFF)
MUTED = RGBColor(0x6E,0x68,0x78)
FAINT = RGBColor(0x9A,0x95,0xA3)
BADBG = RGBColor(0xF8,0xEA,0xEA); BADLN = RGBColor(0xEA, 0xCB, 0xCB)
GOODBG= RGBColor(0xE9,0xF1,0xE9); GOODLN= RGBColor(0xCB, 0xDF, 0xCC)
DCARD = RGBColor(0x2B,0x22,0x3B); DLINE = RGBColor(0x49,0x40,0x59)
DTEXT = RGBColor(0xEC,0xE6,0xDA); DMUTE = RGBColor(0xBE,0xB8,0xC6)

F = "Malgun Gothic"  # Windows 기본 한글 폰트
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill; f.solid(); f.fore_color.rgb = bg
    return s

def _setfont(r, font):
    rPr = r._r.get_or_add_rPr()
    for tag in ('a:latin','a:ea','a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set('typeface', font)

def _spc(p, val):
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    # set char spacing on defRPr
    dr = pPr.find(qn('a:defRPr'))
    if dr is None:
        dr = pPr.makeelement(qn('a:defRPr'), {}); pPr.append(dr)
    dr.set('spc', str(val))

def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf

def line(tf, runs, first=False, align=PP_ALIGN.LEFT, after=6, before=0, lh=1.18, spc=None):
    """runs: list of (text, size, color, bold) or single tuple."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(after); p.space_before = Pt(before)
    p.line_spacing = lh
    if isinstance(runs, tuple): runs = [runs]
    for run in runs:
        text, size, color, bold = (list(run) + [False])[:4]
        font = F
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
        _setfont(r, font)
    if spc is not None: _spc(p, spc)
    return p

def kicker(s, text, y=0.5, dark=False, accentnum=None):
    tf = tb(s, 0.92, y, 11.5, 0.4)
    col = (DMUTE if dark else FAINT)
    p = tf.paragraphs[0]; p.space_after = Pt(0)
    r = p.add_run(); r.text = text.upper()
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = col; r.font.name = MONO
    _setfont(r, MONO); _spc(p, 260)
    return tf

def title(s, text, y=0.95, size=33, dark=False, w=11.5, color=None):
    tf = tb(s, 0.9, y, w, 1.6)
    line(tf, (text, size, color or (DTEXT if dark else INK), True), first=True, lh=1.05, after=0)
    return tf

def accentbar(s, x, y, color=ORANGE, w=0.62, h=0.06):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background(); r.shadow.inherit = False

def card(s, l, t, w, h, fill=PAPER, ln=LINE, lnw=1.0, radius=0.07):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if ln is None: shp.line.fill.background()
    else: shp.line.color.rgb = ln; shp.line.width = Pt(lnw)
    shp.shadow.inherit = False
    try: shp.adjustments[0] = radius
    except Exception: pass
    return shp

def card_text(s, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    tf = tb(s, l+0.18, t+0.14, w-0.36, h-0.28, anchor=anchor)
    for i, ln in enumerate(lines):
        runs = ln['runs']; line(tf, runs, first=(i==0), after=ln.get('after',5),
                                lh=ln.get('lh',1.2), align=ln.get('align',PP_ALIGN.LEFT))
    return tf

def _twidth(text):
    # 한글/넓은 글자는 더 넓게 가중치
    w = 0.0
    for ch in text:
        w += 0.135 if ord(ch) > 0x2000 else 0.075
    return w

def chip(s, l, t, text, dark=False, color=None):
    w = 0.34 + _twidth(text)
    shp = card(s, l, t, w, 0.34, fill=(DCARD if dark else WHITE), ln=(DLINE if dark else LINE), radius=0.5)
    tf = tb(s, l, t, w, 0.34, anchor=MSO_ANCHOR.MIDDLE)
    line(tf, (text, 11, color or (DTEXT if dark else MUTED), True), first=True, align=PP_ALIGN.CENTER, after=0)
    return w

def footer(s, label, idx=None, dark=False):
    num = len(prs.slides._sldIdLst)
    tf = tb(s, 0.9, 7.04, 4, 0.3)
    line(tf, [("■ ", 9, ORANGE, True), ("SIRIAI", 11, (DTEXT if dark else INK), True)], first=True, after=0, spc=120)
    tf2 = tb(s, 9.0, 7.04, 3.43, 0.3)
    line(tf2, ("%02d / 34  ·  %s" % (num, label), 9.5, (DMUTE if dark else FAINT), False), first=True, align=PP_ALIGN.RIGHT, after=0)

def bullets(s, l, t, w, items, dark=False, marker="＋", mcolor=ORANGE, size=14, gap=8, lh=1.3):
    tf = tb(s, l, t, w, 5.2)
    base = DTEXT if dark else INK
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            lead, rest = it
            runs = [(marker+" ", size, mcolor, True), (lead, size, base, True), (rest, size, base, False)]
        else:
            runs = [(marker+" ", size, mcolor, True), (it, size, base, False)]
        line(tf, runs, first=(i==0), after=gap, lh=lh)
    return tf

def ba(s, l, t, w, was, now):
    """before/after mini side by side"""
    half = (w-0.5)/2
    card(s, l, t, half, 0.92, fill=BADBG, ln=BADLN, radius=0.12)
    tf = tb(s, l+0.14, t+0.12, half-0.28, 0.7)
    line(tf, ("기존", 9, RED, True), first=True, after=3, spc=120)
    line(tf, (was, 11.5, INK, False), lh=1.2, after=0)
    tf2 = tb(s, l+half, t, 0.5, 0.92, anchor=MSO_ANCHOR.MIDDLE)
    line(tf2, ("→", 16, FAINT, True), first=True, align=PP_ALIGN.CENTER, after=0)
    card(s, l+half+0.5, t, half, 0.92, fill=GOODBG, ln=GOODLN, radius=0.12)
    tf3 = tb(s, l+half+0.5+0.14, t+0.12, half-0.28, 0.7)
    line(tf3, ("지금", 9, GREEN, True), first=True, after=3, spc=120)
    line(tf3, (now, 11.5, INK, False), lh=1.2, after=0)

def table(s, l, t, w, rows, colw, header_dark=True, fontsizes=(11,11)):
    nrows = len(rows); ncols = len(rows[0])
    gt = s.shapes.add_table(nrows, ncols, Inches(l), Inches(t), Inches(w), Inches(0.4*nrows)).table
    gt.first_row = False; gt.horz_banding = False
    for ci, cw in enumerate(colw): gt.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, cell in enumerate(row):
            c = gt.cell(ri, ci)
            c.margin_left = Inches(0.1); c.margin_right = Inches(0.08)
            c.margin_top = Inches(0.05); c.margin_bottom = Inches(0.05)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.space_after = Pt(0); p.line_spacing = 1.12
            if isinstance(cell, tuple):
                text, color, bold = cell
            else:
                text, color, bold = cell, INK, False
            r = p.add_run(); r.text = text
            sz = fontsizes[0] if ri == 0 else fontsizes[1]
            r.font.size = Pt(sz); r.font.bold = bold or (ri == 0); r.font.name = F
            r.font.color.rgb = (CREAM if ri == 0 and header_dark else color)
            _setfont(r, F)
            cell_fill = c.fill
            if ri == 0:
                cell_fill.solid(); cell_fill.fore_color.rgb = INK
            else:
                cell_fill.solid(); cell_fill.fore_color.rgb = PAPER if ri % 2 else WHITE
    return gt

def dot(s, x, y, color, d=0.12):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = color; o.line.fill.background(); o.shadow.inherit=False

# =========================================================
# SLIDE 1 — COVER
# =========================================================
s = slide(DARK)
kicker(s, "SIRIAI · OPERATIONS PLATFORM · 2026", y=0.7, dark=True)
tf = tb(s, 0.9, 1.5, 11.5, 3.2)
line(tf, ("인플루언서 마케팅,", 50, DTEXT, True), first=True, lh=1.02, after=2)
line(tf, [("이제 ", 50, DTEXT, True), ("시스템", 50, PEACH, True), ("으로", 50, DTEXT, True)], lh=1.02, after=2)
line(tf, ("운영합니다.", 50, DTEXT, True), lh=1.02, after=0)
tf2 = tb(s, 0.92, 4.75, 11, 1.1)
line(tf2, ("1년간 매니저 한 명이 손으로 돌리던 모집·연락·배송·정산·검수를,", 16, DMUTE, False), first=True, lh=1.4, after=2)
line(tf2, [("어드민 웹 · 인플루언서 앱 · AI 자동화", 16, WHITE, True), ("로 다시 설계했습니다.", 16, DMUTE, False)], lh=1.4, after=0)
x = 0.92
for c in ["SIRIAI 어드민", "SIRIAI 인플루언서 앱", "AI 자동화 레이어"]:
    x += chip(s, x, 5.95, c, dark=True, color=DTEXT) + 0.16
tf3 = tb(s, 0.92, 6.7, 11, 0.4)
line(tf3, ("발표자 ____________   ·   SIRIAI   ·   2026.06.07", 12, DMUTE, False), first=True, after=0, spc=120)
footer(s, "표지", 1, dark=True)

# =========================================================
# SLIDE 2 — AGENDA
# =========================================================
s = slide(CREAM)
kicker(s, "CONTENTS", y=0.6)
title(s, "오늘 말씀드릴 순서", y=1.05, size=34)
agenda = [
    ("01", "배경 — 우리가 풀어야 했던 문제", "왜 지금 이걸 만들었나"),
    ("02", "현황 진단 — 5단계에서 무엇이 무너지고 있었나", "모집·연락·배송·정산·검수"),
    ("03", "솔루션 — 두 개의 앱 + AI 자동화", "전체 구조와 문제→해결 매핑"),
    ("04", "단계별 해결 & 실제 화면", "데모 흐름"),
    ("05", "운영 실행 — 어떻게 운영하고 키울 것인가", "연락 · 인바운드 모집 · 인스타 채널"),
    ("06", "효율 · 기술 · 향후 로드맵", "Before/After와 다음 단계"),
]
y = 1.95
for num, t, d in agenda:
    tf = tb(s, 0.95, y, 7.4, 0.62, anchor=MSO_ANCHOR.MIDDLE)
    line(tf, [(num+"   ", 13, FAINT, True), (t, 19, INK, True)], first=True, after=0)
    tf2 = tb(s, 8.4, y, 3.9, 0.62, anchor=MSO_ANCHOR.MIDDLE)
    line(tf2, (d, 11.5, FAINT, False), first=True, align=PP_ALIGN.RIGHT, after=0)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(y+0.72), Inches(11.4), Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background(); ln.shadow.inherit=False
    y += 0.82
footer(s, "목차", 2)

# =========================================================
# SLIDE 3 — 배경
# =========================================================
s = slide(CREAM)
kicker(s, "01 · 배경", y=0.55)
title(s, "규모는 커졌는데, 일하는 방식은 1년 전 그대로였습니다.", y=1.0, size=29, w=11.6)
cards3 = [
    ("🤝  우리가 하는 일", "브랜드와 인플루언서를 잇는 인플루언서 마케팅 대행. 무신사 등 고정 클라이언트부터 신규 대행 캠페인까지."),
    ("📈  무엇이 달라졌나", "고정 클라이언트와 대행 캠페인이 몇 배로 급증. 한 번에 다뤄야 할 인플루언서가 수백 명 규모로."),
    ("🧍  그런데 운영은", "여전히 매니저 1인이 전 과정을 수기로 처리. 검색·문자·시트·송금·다운로드를 사람이 일일이."),
]
cw = 3.73; cx = 0.9; cy = 2.35; ch = 2.5
for i,(h,b) in enumerate(cards3):
    card(s, cx+i*(cw+0.18), cy, cw, ch)
    card_text(s, cx+i*(cw+0.18), cy, cw, ch, [
        {'runs':(h,16,INK,True),'after':8},
        {'runs':(b,13,MUTED,False),'lh':1.45,'after':0},
    ])
card(s, 0.9, 5.2, 11.53, 1.15, fill=BADBG, ln=BADLN, radius=0.09)
tf = tb(s, 1.1, 5.35, 11.1, 0.9, anchor=MSO_ANCHOR.MIDDLE)
line(tf, [("결과   ", 12, RED, True), ("클라이언트 만족도와 회사의 성장성이 동시에 임계점에 도달 — “더 받고 싶어도 더 못 받는” 구조.", 16, INK, True)], first=True, lh=1.3, after=0)
footer(s, "배경", 3)

# =========================================================
# SLIDE 4 — 진단 한눈에
# =========================================================
s = slide(CREAM)
kicker(s, "02 · 현황 진단 · 한눈에", y=0.55)
title(s, "5단계 모든 단계에 수작업 병목이 있었습니다.", y=1.0, size=29, w=11.6)
stages = [("📣","01 모집",BLUE),("💬","02 연락",PURPLE),("📦","03 선정·배송",GREEN),("💸","04 정산",ORANGE),("🎬","05 검수·아카이빙",PINK)]
n=len(stages); gapx=0.2; bw=(11.53-(n-1)*gapx)/n; bx=0.9; by=2.35
for i,(ic,nm,col) in enumerate(stages):
    x=bx+i*(bw+gapx)
    card(s, x, by, bw, 1.5, fill=WHITE, ln=col, lnw=1.5)
    tf=tb(s,x,by+0.18,bw,1.2,anchor=MSO_ANCHOR.TOP)
    line(tf,(ic,24,INK,False),first=True,align=PP_ALIGN.CENTER,after=4)
    line(tf,(nm,13,INK,True),align=PP_ALIGN.CENTER,after=0,lh=1.1)
    if i<n-1:
        ar=tb(s,x+bw-0.02,by,gapx+0.04,1.5,anchor=MSO_ANCHOR.MIDDLE)
        line(ar,("›",18,FAINT,True),first=True,align=PP_ALIGN.CENTER,after=0)
kicker(s, "근본 원인", y=4.3)
roots=[("① 1인 수동 처리 구조","모든 단계를 사람이 직접. 인원이 늘수록 공수 선형 폭증."),
       ("② 데이터 자산화 부재","누가 색조에 강한지, 주소·계좌가 최신인지 — 기억·흩어진 파일 의존."),
       ("③ 공식 시스템·채널 부재","개인 폰·메일·슬랙·시트로 분산. 표준 프로세스 없음.")]
cw=3.73;cy=4.75;ch=1.6
for i,(h,b) in enumerate(roots):
    card(s,0.9+i*(cw+0.18),cy,cw,ch)
    card_text(s,0.9+i*(cw+0.18),cy,cw,ch,[{'runs':(h,14.5,INK,True),'after':6},{'runs':(b,12,MUTED,False),'lh':1.4,'after':0}])
footer(s,"진단 개요",4)

# =========================================================
# PROBLEM SLIDES 5-9
# =========================================================
def problem_slide(idx, num, stage, color, head, items, tail):
    s = slide(CREAM)
    kicker(s, "02 — %s · %s 단계" % (num, stage), y=0.55)
    title(s, head, y=1.0, size=30, w=10.6)
    dot(s, 11.7, 1.12, color, 0.16)
    tf = tb(s, 11.0, 1.05, 1.4, 0.4); line(tf,(stage,12,color,True),first=True,align=PP_ALIGN.RIGHT,after=0)
    accentbar(s, 0.92, 1.95, color)
    # 2x2 panels
    pw=5.65; ph=1.18; px=0.9; py=2.3; gx=0.2; gy=0.18
    for i,it in enumerate(items):
        r=i//2; c=i%2
        x=px+c*(pw+gx); y=py+r*(ph+gy)
        card(s,x,y,pw,ph,fill=BADBG,ln=BADLN,radius=0.1)
        tf=tb(s,x+0.16,y+0.12,pw-0.32,ph-0.24,anchor=MSO_ANCHOR.MIDDLE)
        lead,rest=it
        line(tf,[("✕  ",12,RED,True),(lead,12.5,INK,True),(rest,12.5,INK,False)],first=True,lh=1.28,after=0)
    tf=tb(s,0.9,5.95,11.5,0.7)
    line(tf,[("→ ",14,color,True),(tail[0],14.5,INK,True),(tail[1],14.5,MUTED,False)],first=True,lh=1.3,after=0)
    footer(s,"문제·"+stage,idx)

problem_slide(5,"1","모집",BLUE,"Pool 고갈과 감도 하락",[
    ("아웃바운드(DM)에만 의존 — ","“보내던 사람에게만 또 보내는” 모수 고갈."),
    ("장기 클라이언트 만족도 저하 — ","무신사 등에서 ‘감도 높은 인원 재구성’ 불만 누적."),
    ("크리에이터 성향 자산화 부재 — ","분류가 없어 매번 기억에 의존한 수기 매칭."),
    ("1:1 개별 제안의 리소스 낭비 — ","캠페인마다 일일이 선호·고료를 개별 확인."),
],("신규 캠페인마다 ","태그 검색·수기 1:1 컨택을 처음부터 반복하는 비효율 고착화."))

problem_slide(6,"2","연락",PURPLE,"수작업 커뮤니케이션 과부하",[
    ("공식 채널 부재 · 개인 폰 — ","담당자 개인 휴대폰에 번호를 일일이 저장해 소통."),
    ("다단계 1:1 연락으로 마비 — ","피처체크→선정→가이드→리마인더 3회 전부 수동 문자."),
    ("라포 ↔ 효율의 충돌 — ","단체 시스템이 없어 효율 급락(오픈채팅은 감도와 안 맞음)."),
    ("정보 공개의 딜레마 — ","공공 채널에 올리면 고료가 노출돼 A급 참여 제한."),
],("클라이언트 수가 몇 배로 늘어난 지금, ","담당자 혼자 감당 불가능한 임계점 도달."))

problem_slide(7,"3","배송",GREEN,"데이터 최신화 부재로 인한 배송 사고",[
    ("과거 데이터 재활용 — ","이전 참여자 주소를 확인 없이 그대로 브랜드사에 전달."),
    ("오배송 주기적 발생 — ","이사 등으로 잘못 배송 → 브랜드 신뢰 손상."),
    ("전수 수기 재확인 — ","막으려 한 명 한 명 주소를 다시 묻는 전수조사 추가."),
    ("악순환 — ","사고를 막으려다 공수가 또 늘어남."),
],("결국 ","틀리거나 / 매번 다시 묻거나 — 둘 중 하나였습니다."))

problem_slide(8,"4","정산",ORANGE,"중복 공수 & 패널티 통제 불능",[
    ("중복 서류 제출 피로 — ","여러 번 참여해도 매번 신분증·통장사본 재수취."),
    ("수기 송금의 비효율 — ","메일 취합→슬랙→대표 수기 송금·체크. 누락·지연 위험."),
    ("낮은 고료로 잠수(노쇼) — ","“5만 원 안 받고 말지”로 기한 미준수·연락 두절."),
    ("패널티 기준 부재 — ","늦게 올려도 대부분 정산 → 책임감 결여."),
],("정산이 ","‘돈 주는 일’이 아니라 ‘사람 쫓아다니는 일’이 되어 있었습니다."))

problem_slide(9,"5","검수",PINK,"단순 반복 수작업으로 인한 과부하",[
    ("업로드 링크 수기 입력 — ","수십~수백 명의 링크를 일일이 확인·시트에 복붙."),
    ("전수 다운로드·아카이빙 — ","릴스를 일일이 다운받아 드라이브에 수동 분류·저장."),
    ("시간을 못 쓴 곳 — ","정작 ‘콘텐츠 품질을 보는 일’엔 시간을 못 씀."),
    ("결과 — ","가장 단순한 일에 가장 많은 손이 들어감."),
],("가장 단순한 일에 ","가장 많은 손이 들어가는 구조였습니다."))

# =========================================================
# SLIDE 10 — 전환
# =========================================================
s = slide(DARK)
tf = tb(s, 1.3, 1.7, 10.7, 2.4, anchor=MSO_ANCHOR.MIDDLE)
line(tf, ("“", 60, ORANGE, True), first=True, align=PP_ALIGN.CENTER, after=0, lh=0.5)
line(tf, [("문제는 ", 40, DTEXT, True), ("사람", 40, PEACH, True), ("이 아니라", 40, DTEXT, True)], align=PP_ALIGN.CENTER, after=2, lh=1.1)
line(tf, ("‘시스템의 부재’였습니다.", 40, PEACH, True), align=PP_ALIGN.CENTER, after=0, lh=1.1)
tf2 = tb(s, 1.8, 4.35, 9.7, 1.0)
line(tf2, [("단순 반복은 시스템과 AI가 대신하고, 사람은 ‘감도’와 ‘관계’에 집중하도록 — ", 15, DMUTE, False)], first=True, align=PP_ALIGN.CENTER, lh=1.45, after=2)
line(tf2, ("두 개의 웹앱과 AI 자동화 레이어를 직접 만들었습니다.", 15, WHITE, True), align=PP_ALIGN.CENTER, after=0)
items10 = ["🖥  SIRIAI 어드민 — 운영자", "📱  SIRIAI 인플루언서 앱 — 크리에이터", "🤖  AI 자동화 — 반복 업무"]
total_w = sum(0.34+_twidth(c) for c in items10) + 0.32*2
x = (SW - total_w)/2
for c in items10:
    x += chip(s, x, 5.7, c, dark=True, color=DTEXT) + 0.32
footer(s, "전환", 10, dark=True)

# =========================================================
# SLIDE 11 — 구조도
# =========================================================
s = slide(CREAM)
kicker(s, "03 · 솔루션 · 전체 구조", y=0.55)
title(s, "하나의 데이터베이스를 두 앱이 공유합니다.", y=1.0, size=29)
tf = tb(s, 0.92, 1.75, 11.5, 0.5)
line(tf, [("인플루언서가 앱에서 입력하면 어드민에 ", 14, MUTED, False), ("즉시 반영", 14, INK, True), ("되고, 그 반대도 마찬가지. 다시 묻거나 옮겨 적을 필요가 없습니다.", 14, MUTED, False)], first=True, lh=1.4, after=0)
# three nodes
card(s, 0.9, 2.5, 3.5, 2.3, fill=WHITE)
card_text(s, 0.9, 2.5, 3.5, 2.3, [
    {'runs':("📱  인플루언서 앱",16,INK,True),'after':3,'align':PP_ALIGN.CENTER},
    {'runs':("모바일 · 크리에이터용",11,FAINT,False),'after':8,'align':PP_ALIGN.CENTER},
    {'runs':("· 캠페인 탐색·신청",12,MUTED,False),'after':3,'align':PP_ALIGN.CENTER},
    {'runs':("· 콘텐츠 업로드",12,MUTED,False),'after':3,'align':PP_ALIGN.CENTER},
    {'runs':("· 정산·메시지·알림",12,MUTED,False),'after':0,'align':PP_ALIGN.CENTER},
], anchor=MSO_ANCHOR.MIDDLE)
card(s, 4.92, 2.85, 3.5, 1.6, fill=INK, ln=None)
card_text(s, 4.92, 2.85, 3.5, 1.6, [
    {'runs':("🗄️  공용 DB",16,CREAM,True),'after':4,'align':PP_ALIGN.CENTER},
    {'runs':("Supabase",12,PEACH,True),'after':2,'align':PP_ALIGN.CENTER},
    {'runs':("Postgres · Auth · Storage",11,DMUTE,False),'after':0,'align':PP_ALIGN.CENTER},
], anchor=MSO_ANCHOR.MIDDLE)
for xx in (4.55, 8.5):
    a = tb(s, xx, 2.5, 0.4, 2.3, anchor=MSO_ANCHOR.MIDDLE); line(a, ("⇄", 18, FAINT, True), first=True, align=PP_ALIGN.CENTER, after=0)
card(s, 8.93, 2.5, 3.5, 2.3, fill=WHITE)
card_text(s, 8.93, 2.5, 3.5, 2.3, [
    {'runs':("🖥️  SIRIAI 어드민",16,INK,True),'after':3,'align':PP_ALIGN.CENTER},
    {'runs':("웹 · 운영자용",11,FAINT,False),'after':8,'align':PP_ALIGN.CENTER},
    {'runs':("· 브랜드·캠페인·인플루언서",12,MUTED,False),'after':3,'align':PP_ALIGN.CENTER},
    {'runs':("· 신청·검수·정산",12,MUTED,False),'after':3,'align':PP_ALIGN.CENTER},
    {'runs':("· 콘텐츠 아카이브",12,MUTED,False),'after':0,'align':PP_ALIGN.CENTER},
], anchor=MSO_ANCHOR.MIDDLE)
sub = [("🤖","AI 자동화 레이어"),("🗂️","Google Drive"),("📮","우편번호 API"),("🚀","Vercel 자동배포")]
cw=2.78; cx=0.9; cy=5.05
for i,(ic,t) in enumerate(sub):
    x=cx+i*(cw+0.16)
    card(s,x,cy,cw,1.15)
    tf=tb(s,x,cy+0.16,cw,0.9,anchor=MSO_ANCHOR.TOP)
    line(tf,(ic,18,INK,False),first=True,align=PP_ALIGN.CENTER,after=3)
    line(tf,(t,12.5,INK,True),align=PP_ALIGN.CENTER,after=0)
footer(s,"시스템 구조",11)

# =========================================================
# SLIDE 12 — 문제→해결 매핑
# =========================================================
s = slide(CREAM)
kicker(s, "03 · 문제 → 해결 · 한눈에", y=0.55)
title(s, "5단계의 문제를, 5개의 해결로.", y=1.0, size=29)
rows = [
    ("단계","기존 (문제)","지금 (해결)"),
    (("● 모집",BLUE,True),("DM 의존·기억 매칭으로 Pool 고갈",RED,False),("앱 공개 모집(인바운드) + 구조화 DB·필터 + AI 섭외 스크리닝",INK,True)),
    (("● 연락",PURPLE,True),("개인 폰 다단계 수동 문자·고료 공개 딜레마",RED,False),("자동 알림 + 앱 내 1:1 DM + AI 메시지 초안, 고료는 앱 안에서만",INK,True)),
    (("● 배송",GREEN,True),("과거 주소 재사용→오배송 / 전수 수기 재확인",RED,False),("신청 시마다 주소 직접 입력(우편번호) → 항상 최신 자동 수집",INK,True)),
    (("● 정산",ORANGE,True),("중복 서류·메일→슬랙→수기 송금 누락",RED,False),("계좌 1회 등록 재사용 + 선정 시 자동 생성 + 대시보드·일괄이체",INK,True)),
    (("● 검수",PINK,True),("링크 일일이 복붙·영상 수동 다운로드",RED,False),("인플루언서 직접 업로드 + 드라이브 자동 저장 + AI 검수",INK,True)),
]
table(s, 0.9, 2.05, 11.53, rows, [1.5,4.5,5.53], fontsizes=(11,12))
tf=tb(s,0.9,6.55,11.5,0.5); line(tf,[("→ ",14,ORANGE,True),("다음 장부터, 각 해결을 ",14,MUTED,False),("실제 화면",14,INK,True),("으로 보여드립니다.",14,MUTED,False)],first=True,after=0)
footer(s,"문제→해결 매핑",12)

# =========================================================
# SOLUTION SLIDES 13-17 (text + feature list, right info card)
# =========================================================
def solution_slide(idx, num, stage, color, head, was, now, feats, demo):
    s = slide(CREAM)
    star = "  ★ 핵심" if stage=="정산" else ""
    kicker(s, "04 — %s · 해결 · %s%s" % (num, stage, star), y=0.55)
    title(s, head, y=1.0, size=27, w=10.6)
    dot(s, 11.7, 1.1, color, 0.16)
    tfx=tb(s,11.0,1.03,1.4,0.4); line(tfx,(stage,12,color,True),first=True,align=PP_ALIGN.RIGHT,after=0)
    # before/after mini (left, full width across left col)
    ba(s, 0.9, 1.95, 7.1, was, now)
    bullets(s, 0.9, 3.15, 7.2, feats, size=13.5, gap=9, lh=1.3)
    # right info card
    card(s, 8.35, 1.95, 4.08, 4.55, fill=WHITE)
    tf = tb(s, 8.6, 2.2, 3.6, 4.1)
    line(tf, (demo[0], 12, color, True), first=True, after=10, spc=120)
    for d in demo[1]:
        line(tf, [("▸ ", 12.5, color, True), (d, 12.5, INK, False)], after=8, lh=1.32)
    footer(s, "해결·"+stage, idx)

solution_slide(13,"1","모집",BLUE,"‘찾아 나서는 모집’에서 ‘찾아오게 하는 모집’으로",
    "태그 검색→DM→기억으로 매칭","앱 공개 모집 + 필터로 매칭",
    [("인플루언서 앱 공개 모집 — ","승인 회원이 모집중 캠페인을 직접 보고 ‘신청’. 아웃바운드→인바운드."),
     ("구조화된 인플루언서 DB — ","카테고리·성별·지역·팔로워·협업 단가까지 저장."),
     ("검색·필터 매칭 — ","‘기억’ 대신 ‘조건’으로 추린다."),
     ("AI 섭외 스크리닝 — ","핸들만 넣으면 팔로워·적합도(fit score) 정리.")],
    ("화면 · 인플루언서 앱", ["모집중 캠페인 그리드(2열)","브랜드·고료·콘텐츠 타입 표시","공동작업/2차활용 배지","카테고리·타입 필터"]))

solution_slide(14,"2","연락",PURPLE,"개인 폰을 닫고, 앱 안에서 ‘친절하게 & 한꺼번에’",
    "개인 폰으로 단계마다 수동 문자","앱 자동 알림 + 1:1 DM",
    [("자동 알림 — ","선정🎉/일정📅/심사중⏳/미선정🙏 상태별 자동 표시."),
     ("앱 내 1:1 DM — ","캠페인별 대화 스레드. 개인 폰 없이 라포 유지+효율."),
     ("고료 공개 딜레마 해소 — ","캠페인·고료는 로그인 회원에게만 노출."),
     ("AI 단계별 메시지 — ","섭외·선정·리마인드·정산 메시지 자동 초안.")],
    ("화면 · 알림 & 메시지", ["🎉 ‘문샷 립컬러 선정!’","📅 ‘업로드 일정 06.10~06.17’","⏳ ‘토코보 선크림 심사 중’","💬 관리자 ↔ 인플루언서 1:1"]))

solution_slide(15,"3","배송",GREEN,"주소는 ‘물어보는 것’이 아니라 ‘자동으로 최신화되는 것’",
    "과거 주소 재사용→오배송/전수 재확인","신청 시마다 최신 주소 자동 수집",
    [("신청 = 주소 입력 — ","신청하는 순간 우편번호 검색으로 그때의 주소 직접 입력."),
     ("‘매 캠페인 새로 확인’ 안내 — ","화면에 명시해 오래된 주소가 섞일 여지 제거."),
     ("어드민에 자동 표기 — ","신청자별 배송지가 목록에 바로. 전수조사 소멸."),
     ("브랜드 전달도 그대로 — ","최신 주소를 모아 그대로 넘기면 끝.")],
    ("화면 · 신청 모달", ["다음 우편번호 검색 연동","상세 주소 · 요청사항 입력","‘매 캠페인 새로 확인’ 경고","신청 완료 → 어드민 즉시 반영"]))

solution_slide(16,"4","정산",ORANGE,"‘사람 쫓아다니는 정산’을 ‘자동으로 굴러가는 정산’으로",
    "메일→슬랙→수기 송금·누락","선정 시 자동 생성 + 대시보드",
    [("계좌·주민번호 1회 등록 재사용 — ","분리 보관·서버 조회. 중복 서류 제거."),
     ("선정 → 정산 자동 생성 — ","‘진행’ 선정 시 고료 기본값으로 자동(중복 방지)."),
     ("캠페인별 대시보드 — ","“12명 중 5명 정산완료” 진행바 + 예정 합계."),
     ("1클릭 지급 + AI 보조 — ","지급완료→앱 즉시 반영. 3.3%·일괄이체 CSV.")],
    ("화면 · 정산 대시보드", ["지급 예정 합계 350,000원","📁 문샷: 12명 중 5명 완료","📁 토코보: 8명 전원 완료","선정→자동생성→1클릭 지급"]))

solution_slide(17,"5","검수",PINK,"복붙과 다운로드를, 업로드 한 번과 자동 저장으로",
    "링크 복붙 + 영상 수동 다운로드","직접 업로드 + 드라이브 자동 저장",
    [("인플루언서 직접 업로드 — ","캠페인 선택 후 링크 제출. 담당자 복붙 소멸."),
     ("업로드 전 자가 체크 — ","‘브랜드 태그/공동작업자 추가’ 확인해야 제출."),
     ("어드민 검수 탭 — ","미인증/인증완료 분류, 링크 확인 후 토글."),
     ("아카이브 자동화 + AI 검수 — ","릴스/피드를 드라이브에 자동 저장, 가이드 대조.")],
    ("화면 · 콘텐츠 아카이브", ["릴스 → @핸들.mp4","피드 → @핸들/ 폴더","캠페인·연월별 자동 분류","SIRIAI_콘텐츠 / 문샷 2026.06"]))

# =========================================================
# SLIDE 18 — AI 레이어
# =========================================================
s = slide(CREAM)
kicker(s, "04 — 6 · AI 자동화 레이어", y=0.55)
title(s, "반복되는 ‘쓰기·계산·대조’는 AI가 초안을, 사람은 결정을.", y=1.0, size=27, w=11.6)
ai = [
    ("🔎","influencer-scout","핸들 → 팔로워·카테고리·적합도 스크리닝"),
    ("✉️","influencer-dm","단계별 DM 자동 작성(섭외·리마인드·정산)"),
    ("🧾","campaign-setup","브리프 → 캠페인 등록 + 모집 공고 + 가이드"),
    ("📋","brand-guide","크리에이터 촬영·업로드 가이드 자동 생성"),
    ("✅","content-review","해시태그·멘션·금지사항 가이드 대조 검수"),
    ("💰","settlement-helper","3.3% 원천징수·실지급액·일괄이체 CSV"),
    ("📊","campaign-report","브랜드 보고용 캠페인 성과 리포트"),
    ("📌","ops-digest","리마인드·검수·정산·마감 ‘오늘 할 일’ 점검"),
]
cw=2.78; ch=1.85; gx=0.16; gy=0.18; x0=0.9; y0=2.15
for i,(ic,nm,desc) in enumerate(ai):
    r=i//4; c=i%4
    x=x0+c*(cw+gx); y=y0+r*(ch+gy)
    card(s,x,y,cw,ch)
    card_text(s,x,y,cw,ch,[
        {'runs':(ic,18,INK,False),'after':4},
        {'runs':(nm,13,INK,True),'after':5},
        {'runs':(desc,11,MUTED,False),'lh':1.32,'after':0},
    ])
tf=tb(s,0.9,6.15,11.5,0.6)
line(tf,[("→ 캠페인 셋업부터 섭외·메시지·검수·정산·보고까지 — ",13.5,MUTED,False),("라이프사이클 전체에 자동화 도우미.",13.5,INK,True)],first=True,after=0)
footer(s,"AI 자동화",18)

# =========================================================
# SLIDE — 운영 실행 · 개요
# =========================================================
s = slide(CREAM)
kicker(s, "05 · 운영 실행 · 개요", y=0.55)
title(s, "앱을 만든 게 끝이 아닙니다 — 이렇게 운영하고 키웁니다.", y=1.0, size=26, w=11.9)
tf = tb(s, 0.92, 1.74, 11.6, 0.55)
line(tf, [("회사에서 가장 궁금해할 부분 — 시스템을 ", 13.5, MUTED, False), ("실제로 어떻게 굴리고", 13.5, INK, True), (", 신규 인플루언서를 ", 13.5, MUTED, False), ("어떻게 모으고", 13.5, INK, True), (", 채널을 ", 13.5, MUTED, False), ("어떻게 키울지", 13.5, INK, True), (".", 13.5, MUTED, False)], first=True, lh=1.4, after=0)
ax = [("축 1", PURPLE, "대량 연락", "개인 폰 대신 문자 API·카카오 알림톡으로, 앱 DB와 연동해 단계별 자동·일괄 발송."),
      ("축 2", BLUE, "인바운드 모집", "인스타 허브 + 레퍼럴 + 기존 DB로 ‘찾아오게’ 만드는 신규 유입 엔진."),
      ("축 3", GREEN, "인스타 채널", "저장하고 싶은 콘텐츠로 브랜딩 → 앱 유입·락인 → 다시 모집으로 순환.")]
cw = 3.73; cy = 2.5; ch = 2.0
for i,(tg,col,h,d) in enumerate(ax):
    x = 0.9 + i*(cw+0.18)
    card(s, x, cy, cw, ch)
    pw = 0.3 + _twidth(tg)
    pp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.18), Inches(cy+0.18), Inches(pw), Inches(0.32))
    pp.fill.solid(); pp.fill.fore_color.rgb = col; pp.line.fill.background(); pp.shadow.inherit = False
    try: pp.adjustments[0] = 0.5
    except Exception: pass
    ptf = pp.text_frame; ptf.margin_top=0; ptf.margin_bottom=0; pq = ptf.paragraphs[0]; pq.alignment = PP_ALIGN.CENTER
    rr = pq.add_run(); rr.text = tg; rr.font.size = Pt(10.5); rr.font.bold = True; rr.font.color.rgb = WHITE; rr.font.name = F; _setfont(rr, F)
    tf = tb(s, x+0.18, cy+0.64, cw-0.36, ch-0.78)
    line(tf, (h, 15, INK, True), first=True, after=6)
    line(tf, (d, 11.5, MUTED, False), lh=1.35, after=0)
kicker(s, "운영 원칙 3", y=4.78)
pr = [("① 단순 반복은 자동화", "문자·검수·정산·콘텐츠 초안은 시스템·AI가. 사람은 판단·관계."),
      ("② 아웃바운드 → 인바운드", "찾아 나서는 모집을, 찾아오는 구조로 전환."),
      ("③ 한 번 만들어 여러 번", "결과물·콘텐츠를 재활용(아카이브·크로스미디어).")]
cy2 = 5.22; ch2 = 1.4
for i,(h,d) in enumerate(pr):
    x = 0.9 + i*(cw+0.18)
    card(s, x, cy2, cw, ch2)
    card_text(s, x, cy2, cw, ch2, [{'runs':(h,13,INK,True),'after':5}, {'runs':(d,11,MUTED,False),'lh':1.3,'after':0}])
footer(s, "운영·개요")

# =========================================================
# SLIDE 19 — 운영 ① 대량 연락
# =========================================================
s = slide(CREAM)
kicker(s, "05 · 운영 실행 · 대량 연락", y=0.55)
title(s, "개인 폰 없이, 수백 명에게 한 번에.", y=1.0, size=29)
tf = tb(s, 0.92, 1.78, 11.5, 0.6)
line(tf, [("인플루언서 앱 DB(상태·단계)와 연동해 ", 14, MUTED, False), ("맞춤 대량 발송", 14, INK, True), (". 선정·리마인드·정산 안내를 단계에 맞춰 버튼 한 번으로.", 14, MUTED, False)], first=True, lh=1.4, after=0)
cw = 5.66; cy = 2.5; ch = 2.6
card(s, 0.9, cy, cw, ch)
card_text(s, 0.9, cy, cw, ch, [
    {'runs':("📨  방법 1 · 단체 문자 발송 API", 15.5, INK, True), 'after':10},
    {'runs':[("· ", 12.5, ORANGE, True), ("알리고 · 문자나라 · 솔라피 등 문자 API 연동", 12.5, MUTED, False)], 'after':7, 'lh':1.3},
    {'runs':[("· ", 12.5, ORANGE, True), ("한 번에 수십~수백 명, 이름·캠페인 자동 치환 맞춤 문자", 12.5, MUTED, False)], 'after':7, 'lh':1.3},
    {'runs':[("· ", 12.5, ORANGE, True), ("비용 월 몇만 원 수준", 12.5, MUTED, False)], 'after':0, 'lh':1.3},
])
card(s, 0.9+cw+0.2, cy, cw, ch)
card_text(s, 0.9+cw+0.2, cy, cw, ch, [
    {'runs':("💬  방법 2 · 카카오 알림톡", 15.5, INK, True), 'after':10},
    {'runs':[("· ", 12.5, ORANGE, True), ("카카오 비즈니스 채널 개설 후 발송", 12.5, MUTED, False)], 'after':7, 'lh':1.3},
    {'runs':[("· ", 12.5, ORANGE, True), ("카톡 형식이라 친숙 · 개인 번호 노출 없이 대량", 12.5, MUTED, False)], 'after':7, 'lh':1.3},
    {'runs':[("· ", 12.5, ORANGE, True), ("읽힘률이 일반 문자보다 높음", 12.5, MUTED, False)], 'after':0, 'lh':1.3},
])
card(s, 0.9, 5.55, 11.52, 0.92, fill=RGBColor(0xEA,0xF1,0xFB), ln=RGBColor(0xCB,0xDD,0xF5), radius=0.12)
tf = tb(s, 1.1, 5.62, 11.1, 0.78, anchor=MSO_ANCHOR.MIDDLE)
line(tf, [("현재 앱 내 알림 → 문자·알림톡으로 휴대폰까지 도달 강화. ", 14, INK, True), ("둘 다 앱과 연동한 단체 발송이 핵심.", 14, MUTED, False)], first=True, lh=1.3, after=0)
footer(s, "운영·연락", 19)

# =========================================================
# SLIDE — 운영·연락 단계별 시나리오
# =========================================================
s = slide(CREAM)
kicker(s, "05 · 운영 실행 · 단계별 자동 발송", y=0.55)
title(s, "캠페인 한 건의 연락, 이렇게 자동화됩니다.", y=1.0, size=28)
tf = tb(s, 0.92, 1.74, 11.6, 0.55)
line(tf, [("예전엔 캠페인당 [피처체크→선정→가이드→리마인더 3회→정산]을 전부 손으로 문자. 이제는 앱 상태에 맞춰 ", 13.5, MUTED, False), ("템플릿 + 버튼 한 번", 13.5, INK, True), ("으로.", 13.5, MUTED, False)], first=True, lh=1.4, after=0)
rows = [("단계", "발송 내용", "채널", "자동화 포인트"),
    (("① 섭외·모집 안내", INK, True), ("캠페인 오픈 · 신청 유도", INK, False), ("알림톡 · 문자", MUTED, False), ("DB 필터 대상에 일괄", GREEN, True)),
    (("② 선정 안내", INK, True), ("“선정되셨어요” + 일정", INK, False), ("앱 알림 + 알림톡", MUTED, False), ("‘진행’ 선정 시 자동", GREEN, True)),
    (("③ 가이드 전달", INK, True), ("촬영·업로드 가이드 링크", INK, False), ("앱 + 알림톡", MUTED, False), ("캠페인 가이드 자동 첨부", GREEN, True)),
    (("④ 리마인드 ×3", INK, True), ("3일 전 · 1일 전 · 당일", INK, False), ("알림톡 · 문자", MUTED, False), ("업로드 마감 기준 예약", GREEN, True)),
    (("⑤ 정산 안내", INK, True), ("지급 예정 · 완료 통지", INK, False), ("앱 + 알림톡", MUTED, False), ("지급완료 시 자동", GREEN, True)),
]
table(s, 0.9, 2.45, 11.53, rows, [2.4, 3.6, 2.9, 2.63], fontsizes=(11, 11.5))
tf = tb(s, 0.9, 6.45, 11.5, 0.5)
line(tf, [("→ AI ", 13.5, MUTED, False), ("influencer-dm", 13.5, INK, True), (" 스킬이 각 단계 메시지 초안을 작성, 담당자는 확인·발송만.", 13.5, MUTED, False)], first=True, after=0)
footer(s, "운영·연락 시나리오")

# =========================================================
# SLIDE 20 — 운영 ② 인바운드 모집 엔진
# =========================================================
s = slide(CREAM)
kicker(s, "05 · 운영 실행 · 인바운드 모집 · 개요", y=0.55)
title(s, "찾아 나서지 않고, 찾아오게 만든다.", y=1.0, size=29)
eng = [
    ("🏢","인스타 = 인바운드 허브","“감도 높은 뷰티·패션 브랜드 × 크리에이터” 포지셔닝. 프로필에 앱 링크."),
    ("📲","상시 노출","진행 결과물(공동작업자 릴스·피드) 노출 + 스토리에 ‘모집중 캠페인’ 정기 업로드."),
    ("🤝","레퍼럴(지인 추천)","기존 인플루언서에게 지인 추천 요청. 진행 성사 시 소개 리워드."),
    ("🔔","기존 DB 활성화","기존 협업 인플루언서에게 메일·문자·알림톡으로 앱 가입 유도."),
]
cw = 2.78; cy = 2.2; ch = 2.35
for i,(ic,h,d) in enumerate(eng):
    x = 0.9 + i*(cw+0.16)
    card(s, x, cy, cw, ch)
    card_text(s, x, cy, cw, ch, [
        {'runs':(ic, 18, INK, False), 'after':4},
        {'runs':(h, 12.5, INK, True), 'after':6, 'lh':1.15},
        {'runs':(d, 10.5, MUTED, False), 'lh':1.32, 'after':0},
    ])
kicker(s, "성장 플라이휠", y=4.95)
fly = ["① 인스타 노출·결과물","② 앱 가입·신청","③ 캠페인 진행","④ 고퀄 결과물","↻ 다시 콘텐츠+레퍼럴"]
fcw = 2.1; fgap = (11.52 - fcw*5)/4; fx = 0.9; fy = 5.42
for i,t in enumerate(fly):
    x = fx + i*(fcw+fgap)
    last = (i==4)
    card(s, x, fy, fcw, 0.6, fill=(RGBColor(0xEA,0xF1,0xFB) if last else WHITE), ln=(BLUE if last else LINE), radius=0.2)
    tf = tb(s, x+0.04, fy, fcw-0.08, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    line(tf, (t, 10, (BLUE if last else INK), True), first=True, align=PP_ALIGN.CENTER, after=0, lh=1.0)
    if i < 4:
        a = tb(s, x+fcw-0.03, fy, fgap+0.06, 0.6, anchor=MSO_ANCHOR.MIDDLE)
        line(a, ("→", 13, FAINT, True), first=True, align=PP_ALIGN.CENTER, after=0)
tf = tb(s, 0.9, 6.28, 11.5, 0.5)
line(tf, [("→ 광고비에 기대지 않고, ", 13.5, MUTED, False), ("결과물이 다음 모집을 부르는 선순환.", 13.5, INK, True)], first=True, after=0)
footer(s, "운영·모집", 20)

# =========================================================
# SLIDE — 모집 ① 인스타 허브
# =========================================================
s = slide(CREAM)
kicker(s, "05 · 운영 실행 · 모집 ① 인스타 허브", y=0.55)
title(s, "회사 인스타그램 = 인바운드의 중심.", y=1.0, size=29)
bullets(s, 0.9, 2.1, 6.9, [
    ("포지셔닝 — ", "“감도 높은 뷰티·패션 브랜드와 크리에이터를 잇는다”를 한 줄로 각인."),
    ("프로필 = 전환 장치 — ", "소개문 + 링크바이오에 인플루언서 앱 링크(가입·모집중 캠페인) 고정."),
    ("운영 목표 — ", "팔로워(=잠재 인플루언서 풀)를 모으고 자연스럽게 앱으로 유입."),
    ("두 대상 동시 공략 — ", "인플루언서엔 ‘여기서 협업하고 싶다’, 브랜드엔 ‘대행 맡기고 싶다’."),
], size=13.5, gap=11)
card(s, 8.1, 2.1, 4.33, 2.95)
tf = tb(s, 8.35, 2.32, 3.85, 2.6)
line(tf, ("프로필 예시", 10.5, FAINT, True), first=True, after=9, spc=160)
line(tf, ("SIRIAI · 뷰티·패션 인플루언서 마케팅", 13, INK, True), after=6, lh=1.4)
line(tf, ("─ 감도 높은 브랜드 × 크리에이터 매칭", 12, MUTED, False), after=5, lh=1.4)
line(tf, ("─ 진행 캠페인 포트폴리오 큐레이션", 12, MUTED, False), after=5, lh=1.4)
line(tf, ("─ 지금 모집중 캠페인 신청 👇", 12, MUTED, False), after=5, lh=1.4)
line(tf, ("🔗 link.in/siriai-app", 12, BLUE, True), after=0)
card(s, 0.9, 5.35, 11.52, 0.92, fill=GOODBG, ln=GOODLN, radius=0.12)
tf = tb(s, 1.1, 5.45, 11.1, 0.72, anchor=MSO_ANCHOR.MIDDLE)
line(tf, [("핵심 — 인스타는 ‘광고판’이 아니라 앱으로 가는 입구. ", 14, INK, True), ("모든 콘텐츠 끝에 ‘앱에서 신청·확인’ 동선.", 14, MUTED, False)], first=True, lh=1.3, after=0)
footer(s, "모집·인스타 허브")

# =========================================================
# SLIDE — 모집 ② 노출 전략
# =========================================================
s = slide(CREAM)
kicker(s, "05 · 운영 실행 · 모집 ② 노출 전략", y=0.55)
title(s, "결과물 자체가 광고가 되게 한다.", y=1.0, size=29)
ex = [("🤝", "공동작업자(협업) 노출", "진행 인플루언서 게시물에 회사 계정을 공동작업자로 → 그들의 피드에 함께 노출. 그 게시물로 광고 집행까지 가능."),
      ("🎞️", "결과물 재업로드", "진행한 릴스·피드를 인플루언서 태그 포함해 회사 계정에 재업로드 → 포트폴리오 + 자연 도달."),
      ("📲", "스토리 상시 모집", "“현재 모집중 캠페인” 스토리를 정기 업로드 → 팔로우 중인 기존 인플루언서가 보고 바로 신청.")]
cw = 3.73; cy = 2.3; ch = 2.95
for i,(ic,h,d) in enumerate(ex):
    x = 0.9 + i*(cw+0.18)
    card(s, x, cy, cw, ch)
    card_text(s, x, cy, cw, ch, [{'runs':(ic,20,INK,False),'after':6}, {'runs':(h,14,INK,True),'after':7,'lh':1.2}, {'runs':(d,11.5,MUTED,False),'lh':1.4,'after':0}])
tf = tb(s, 0.9, 5.6, 11.5, 0.6)
line(tf, [("→ 별도 제작물 없이 ", 14, MUTED, False), ("이미 가진 캠페인 결과물", 14, INK, True), ("로 노출을 만들고, 그 노출이 다시 신청으로.", 14, MUTED, False)], first=True, after=0)
footer(s, "모집·노출 전략")

# =========================================================
# SLIDE — 모집 ③ 레퍼럴 · 기존 DB
# =========================================================
s = slide(CREAM)
kicker(s, "05 · 운영 실행 · 모집 ③ 레퍼럴 · 기존 DB", y=0.55)
title(s, "검증된 사람이, 검증된 사람을 데려온다.", y=1.0, size=29)
cw = 5.66; cy = 2.3; ch = 3.0
card(s, 0.9, cy, cw, ch, fill=GOODBG, ln=GOODLN)
card_text(s, 0.9, cy, cw, ch, [
    {'runs':("🤝  레퍼럴 (지인 추천)", 14, GREEN, True), 'after':11},
    {'runs':[("✓  ",12,GREEN,True),("진행 이력 있는 인플루언서에게 ‘비슷한 감도의 지인 추천’ 요청.",12.5,INK,False)], 'after':9, 'lh':1.35},
    {'runs':[("✓  ",12,GREEN,True),("소개 인센티브 — 가입·승인 + 실제 캠페인 진행 시 소정의 리워드.",12.5,INK,False)], 'after':9, 'lh':1.35},
    {'runs':[("✓  ",12,GREEN,True),("‘성과 연동’이라 비용 효율적 + 감도 검증된 풀 유입.",12.5,INK,False)], 'after':0, 'lh':1.35},
])
card(s, 0.9+cw+0.2, cy, cw, ch, fill=GOODBG, ln=GOODLN)
card_text(s, 0.9+cw+0.2, cy, cw, ch, [
    {'runs':("🔔  기존 DB 활성화", 14, GREEN, True), 'after':11},
    {'runs':[("✓  ",12,GREEN,True),("그동안 협업한 인플루언서 전체에게 앱 가입 안내(메일·문자·카톡).",12.5,INK,False)], 'after':9, 'lh':1.35},
    {'runs':[("✓  ",12,GREEN,True),("“이제 신청·정산·소통이 앱 한 곳에서” 가치를 전달.",12.5,INK,False)], 'after':9, 'lh':1.35},
    {'runs':[("✓  ",12,GREEN,True),("잠들어 있던 풀을 활성 사용자로 전환.",12.5,INK,False)], 'after':0, 'lh':1.35},
])
tf = tb(s, 0.9, 5.55, 11.5, 0.6)
line(tf, [("→ 광고 의존 없이 ", 14, MUTED, False), ("신뢰 기반", 14, INK, True), ("으로 풀을 늘리고, 자발적 추천으로 지속 유입.", 14, MUTED, False)], first=True, after=0)
footer(s, "모집·레퍼럴·DB")

# =========================================================
# SLIDE 21 — 운영 ③ 인스타 콘텐츠 전략
# =========================================================
s = slide(CREAM)
kicker(s, "05 · 운영 실행 · 인스타 채널 · 개요", y=0.5)
title(s, "팔리는 광고가 아니라, ‘저장하고 싶은 채널’.", y=0.92, size=27)
kicker(s, "콘텐츠 4기둥", y=1.66)
pil = [
    ("포트폴리오",BLUE,"크리에이터 아카이브 — 진행 결과물 큐레이션. ‘여기서 하면 퀄리티 브랜드를 만난다’ 증명 + ‘나도 협찬받고 싶다’ 자극."),
    ("트렌드·매거진",GREEN,"분기 뷰티·패션 트렌드 리포트, ‘성수 라이징 브랜드 5’ 류 카드뉴스. 마지막 장 ‘친구 공유’로 바이럴."),
    ("교육형",PURPLE,"3초 후킹·GRWM·조명·자막·‘잘 되는 캡션 공식’·트렌드 사운드. 크리에이터 성장 지원으로 락인."),
    ("모집·전환",ORANGE,"‘신청하면 이렇게 진행돼요’ 3단계, 모집중 캠페인 티저, 레퍼럴 이벤트."),
]
cw = 2.78; cy = 2.0; ch = 2.45
for i,(tg,col,d) in enumerate(pil):
    x = 0.9 + i*(cw+0.16)
    card(s, x, cy, cw, ch)
    pw = 0.3 + _twidth(tg)
    pp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.18), Inches(cy+0.18), Inches(pw), Inches(0.34))
    pp.fill.solid(); pp.fill.fore_color.rgb = col; pp.line.fill.background(); pp.shadow.inherit = False
    try: pp.adjustments[0] = 0.5
    except Exception: pass
    ptf = pp.text_frame; ptf.margin_top=0; ptf.margin_bottom=0; pq = ptf.paragraphs[0]; pq.alignment = PP_ALIGN.CENTER
    rr = pq.add_run(); rr.text = tg; rr.font.size = Pt(11); rr.font.bold = True; rr.font.color.rgb = WHITE; rr.font.name = F; _setfont(rr, F)
    tf = tb(s, x+0.18, cy+0.68, cw-0.36, ch-0.82)
    line(tf, (d, 11.5, MUTED, False), first=True, lh=1.32, after=0)
strip = [
    ("🔁 크로스미디어","인스타엔 핵심 요약 → ‘전체는 앱에서’ CTA → 앱 유입. 앱 내 뉴스레터로 상시 열람·락인."),
    ("📢 소액 광고","뷰티·패션 타깃 일 1~2만 원 ‘협업 모집’ 광고 + 공동작업자 노출로 도달 확대."),
    ("🤖 올인원 자동화","AI가 주제 선정·트렌드 데이터 수집·카드뉴스 이미지+캡션까지 자동 생성."),
]
cw2 = 3.73; sy = 4.7; sh2 = 1.7
for i,(h,d) in enumerate(strip):
    x = 0.9 + i*(cw2+0.18)
    card(s, x, sy, cw2, sh2)
    card_text(s, x, sy, cw2, sh2, [
        {'runs':(h, 13, INK, True), 'after':5},
        {'runs':(d, 11, MUTED, False), 'lh':1.32, 'after':0},
    ])
footer(s, "운영·콘텐츠", 21)

# =========================================================
# SLIDE — 콘텐츠 기둥 ① 포트폴리오
# =========================================================
s = slide(CREAM)
kicker(s, "05 · 운영 실행 · 콘텐츠 기둥 ① 포트폴리오", y=0.55)
title(s, "“여기서 하면 퀄리티 브랜드를 만난다”를 증명.", y=1.0, size=26, w=11.9)
bullets(s, 0.9, 2.1, 6.9, [
    ("크리에이터 아카이브 — ", "매칭된 인플루언서들이 촬영한 고퀄 피드·릴스를 큐레이션해 업로드."),
    ("포트폴리오 증명 — ", "‘이 브랜드의 이 제품을 이렇게 소화’ → 브랜드사엔 대행 능력 증명."),
    ("욕구 자극 — ", "다른 인플루언서엔 ‘나도 저기서 협찬받아 예쁜 피드를 만들고 싶다’."),
    ("비주얼 브랜드 소개 — ", "진행 브랜드(때론 하이엔드 룩북)도 감성 이미지로 소개해 계정 격 상승."),
], size=13, gap=11)
card(s, 8.1, 2.1, 4.33, 3.5)
tf = tb(s, 8.35, 2.32, 3.85, 3.1)
line(tf, ("게시물 예시", 10.5, FAINT, True), first=True, after=10, spc=160)
for t in ["“[브랜드] 신상 립 — 우리 크리에이터들의 해석”", "릴스 모음: 이번 달 베스트 콘텐츠 5", "비포&애프터: 같은 제품, 다른 무드", "협업 브랜드 신제품·행사 소개"]:
    line(tf, [("· ", 12, ORANGE, True), (t, 12, MUTED, False)], after=9, lh=1.35)
footer(s, "콘텐츠·포트폴리오")

# =========================================================
# SLIDE — 콘텐츠 기둥 ② 트렌드·매거진
# =========================================================
s = slide(CREAM)
kicker(s, "05 · 운영 실행 · 콘텐츠 기둥 ② 트렌드·매거진", y=0.55)
title(s, "‘저장’과 ‘공유’를 부르는 정보성 콘텐츠.", y=1.0, size=28)
tf = tb(s, 0.92, 1.74, 11.6, 0.6)
line(tf, [("광고만 올리면 매력도가 떨어집니다. 저장하고 싶은 매거진형 콘텐츠로 도달·팔로워를 키웁니다. (기획 중인 ", 13.5, MUTED, False), ("자체 뷰티 트렌드 리포트", 13.5, INK, True), ("와 직결)", 13.5, MUTED, False)], first=True, lh=1.4, after=0)
cw = 5.66; cy = 2.55; ch = 3.05
card(s, 0.9, cy, cw, ch)
card_text(s, 0.9, cy, cw, ch, [
    {'runs':("주제 예시", 14, INK, True), 'after':10},
    {'runs':[("· ",12,ORANGE,True),("“올여름 파우치 치트키, 젤리 립 플럼퍼 Top 5”",12,MUTED,False)], 'after':9, 'lh':1.3},
    {'runs':[("· ",12,ORANGE,True),("“요즘 성수에서 가장 핫한 라이징 브랜드 5”",12,MUTED,False)], 'after':9, 'lh':1.3},
    {'runs':[("· ",12,ORANGE,True),("“올여름 주목할 뷰티 트렌드 키워드”",12,MUTED,False)], 'after':9, 'lh':1.3},
    {'runs':[("· ",12,ORANGE,True),("“브랜드가 요즘 찾는 인플루언서 유형”",12,MUTED,False)], 'after':0, 'lh':1.3},
])
card(s, 0.9+cw+0.2, cy, cw, ch)
card_text(s, 0.9+cw+0.2, cy, cw, ch, [
    {'runs':("운영 포인트", 14, INK, True), 'after':10},
    {'runs':[("· ",12,ORANGE,True),("분기별 트렌드 리포트를 카드뉴스로 제작·발행",12,MUTED,False)], 'after':9, 'lh':1.3},
    {'runs':[("· ",12,ORANGE,True),("시즌·테마별 큐레이션으로 ‘저장’ 유도",12,MUTED,False)], 'after':9, 'lh':1.3},
    {'runs':[("· ",12,ORANGE,True),("마지막 장 ‘친구에게 공유’ → 바이럴",12,MUTED,False)], 'after':9, 'lh':1.3},
    {'runs':[("· ",12,ORANGE,True),("앱 내 뉴스레터로도 상시 열람 제공",12,MUTED,False)], 'after':0, 'lh':1.3},
])
footer(s, "콘텐츠·트렌드")

# =========================================================
# SLIDE — 콘텐츠 기둥 ③ 교육 · ④ 모집
# =========================================================
s = slide(CREAM)
kicker(s, "05 · 운영 실행 · 콘텐츠 기둥 ③ 교육 · ④ 모집", y=0.55)
title(s, "크리에이터를 키우고(락인), 신청으로 전환한다.", y=1.0, size=26, w=11.9)
def _tagpill(x, y, tg, col):
    pw = 0.3 + _twidth(tg)
    pp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(pw), Inches(0.32))
    pp.fill.solid(); pp.fill.fore_color.rgb = col; pp.line.fill.background(); pp.shadow.inherit = False
    try: pp.adjustments[0] = 0.5
    except Exception: pass
    t = pp.text_frame; t.margin_top=0; t.margin_bottom=0; q = t.paragraphs[0]; q.alignment = PP_ALIGN.CENTER
    r = q.add_run(); r.text = tg; r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = F; _setfont(r, F)
cw = 5.66; cy = 2.2; ch = 3.45
card(s, 0.9, cy, cw, ch)
_tagpill(0.9+0.18, cy+0.18, "기둥 ③ 교육형", PURPLE)
tf = tb(s, 0.9+0.18, cy+0.64, cw-0.36, ch-0.78)
for j,t in enumerate(["“3초 후킹 만드는 법” · GRWM 촬영 가이드", "조명·각도 꿀팁 · 무료 편집앱·템플릿", "자막 스타일 · ‘잘 되는 캡션 공식’", "트렌드 사운드 · 알고리즘·도달 올리는 법", "선정 확률 높이는 프로필 관리법(=앱 활용)"]):
    line(tf, [("· ",12,PURPLE,True),(t,12,MUTED,False)], first=(j==0), after=7, lh=1.3)
card(s, 0.9+cw+0.2, cy, cw, ch)
_tagpill(0.9+cw+0.2+0.18, cy+0.18, "기둥 ④ 모집·전환", ORANGE)
tf = tb(s, 0.9+cw+0.2+0.18, cy+0.64, cw-0.36, ch-0.78)
for j,t in enumerate(["“신청하면 이렇게 진행돼요” 3단계 안내", "한정 캠페인 티저 · 마감 임박 리마인드", "레퍼럴 이벤트(‘좋은 동료를 추천해주세요’)", "앱 활용 꿀팁 + ‘자세한 건 앱에서’ CTA"]):
    line(tf, [("· ",12,ORANGE,True),(t,12,MUTED,False)], first=(j==0), after=7, lh=1.3)
footer(s, "콘텐츠·교육·모집")

# =========================================================
# SLIDE — 콘텐츠 크로스미디어·자동화
# =========================================================
s = slide(CREAM)
kicker(s, "05 · 운영 실행 · 크로스미디어 · 자동화", y=0.55)
title(s, "인스타 → 앱 유입, 그리고 운영은 자동화.", y=1.0, size=28)
cm = [("🔁", "크로스미디어 전략", "같은 가이드를 인스타용 카드뉴스로 요약 → ‘전체는 앱에서’ 링크 → 팔로워를 앱으로 유입."),
      ("📱", "앱 락인", "성장 가이드·콘텐츠 팁·뉴스레터를 앱 내 상시 제공 → 다운로드·재방문 유도."),
      ("📢", "소액 광고", "뷰티·패션 크리에이터 타깃 ‘협업 모집’ 광고를 일 1~2만 원으로 → 신규 인바운드.")]
cw = 3.73; cy = 2.2; ch = 2.2
for i,(ic,h,d) in enumerate(cm):
    x = 0.9 + i*(cw+0.18)
    card(s, x, cy, cw, ch)
    card_text(s, x, cy, cw, ch, [{'runs':(ic,18,INK,False),'after':5}, {'runs':(h,13.5,INK,True),'after':6}, {'runs':(d,11.5,MUTED,False),'lh':1.35,'after':0}])
kicker(s, "올인원 자동화 시스템", y=4.78)
au = [("① 기획", "업로드할 최적의 주제·아이템을 AI가 선정."), ("② 수집", "트렌드 분석용 데이터를 시스템이 자동 수집."), ("③ 제작", "카드뉴스 이미지 + 본문 캡션까지 자동 생성.")]
cy2 = 5.22; ch2 = 1.42
for i,(h,d) in enumerate(au):
    x = 0.9 + i*(cw+0.18)
    card(s, x, cy2, cw, ch2)
    card_text(s, x, cy2, cw, ch2, [{'runs':(h,13,INK,True),'after':5}, {'runs':(d,11,MUTED,False),'lh':1.3,'after':0}])
footer(s, "콘텐츠·자동화")

# =========================================================
# SLIDE 22 — Before/After
# =========================================================
s = slide(CREAM)
kicker(s, "06 · 업무 효율 · Before / After", y=0.55)
title(s, "같은 일을, 더 적은 손으로.", y=1.0, size=29)
rows19 = [
    ("항목","Before · 수작업","After · 시스템 + AI"),
    ("모집",("태그 검색 + 1:1 DM, 기억 매칭",RED,False),("앱 공개 모집 + 조건 필터 + AI 스카우트",INK,True)),
    ("소통",("개인 폰 수동 문자 × 인원수",RED,False),("자동 알림 + 앱 1:1 DM (+AI 초안)",INK,True)),
    ("배송지",("전수 수기 재확인 / 과거 주소 사고",RED,False),("신청 시 최신 주소 자동 수집",INK,True)),
    ("서류",("매 캠페인 신분증·통장 재수취",RED,False),("1회 등록 후 재사용",INK,True)),
    ("송금",("메일 취합→슬랙→수기 송금",RED,False),("선정 시 자동 생성 + 대시보드 + 일괄이체",INK,True)),
    ("검수·저장",("링크 복붙 + 영상 수동 다운로드",RED,False),("직접 업로드 + 드라이브 자동 + AI 검수",INK,True)),
    ("운영 주체",("1인이 전 과정 수작업",RED,False),("시스템·AI가 반복, 사람은 판단·관계",INK,True)),
]
table(s, 0.9, 1.85, 11.53, rows19, [1.7,4.7,5.13], fontsizes=(11,11.5))
card(s, 0.9, 6.35, 11.53, 0.75, fill=GOODBG, ln=GOODLN, radius=0.12)
tf=tb(s,1.1,6.45,11.1,0.55,anchor=MSO_ANCHOR.MIDDLE)
line(tf,[("핵심 — ",14,GREEN,True),("인원이 늘어도 매니저 공수는 선형으로 늘지 않는다. 규모를 키울 수 있는 구조.",14,INK,True)],first=True,after=0)
footer(s,"Before / After",22)

# =========================================================
# SLIDE 20 — 기술·보안·운영
# =========================================================
s = slide(CREAM)
kicker(s, "06 · 기술 · 보안 · 운영", y=0.55)
title(s, "검증된 표준 스택 위에, 안전하게.", y=1.0, size=29)
cols20 = [
    ("🧱  기술 스택",["Next.js + TypeScript","Supabase (DB·인증·스토리지)","Vercel 자동 배포","Google Drive · 우편번호 API"]),
    ("🔒  보안 · 개인정보",["주민번호 별도 테이블 분리 보관","서버에서 권한 키로만 조회","가입은 검토 후 승인제","정산 정보는 본인·담당자만"]),
    ("⚙️  운영",["코드 반영 → 즉시 자동 배포","개발 현황 페이지로 진척 관리","웹앱이라 설치 없이 링크 접속","두 앱 + AI를 자체 운영"]),
]
cw=3.73; cy=2.2; ch=3.4
for i,(h,lst) in enumerate(cols20):
    x=0.9+i*(cw+0.18)
    card(s,x,cy,cw,ch)
    tf=tb(s,x+0.2,cy+0.2,cw-0.4,ch-0.4)
    line(tf,(h,15.5,INK,True),first=True,after=10)
    for it in lst:
        line(tf,[("· ",12.5,ORANGE,True),(it,12.5,MUTED,False)],after=7,lh=1.3)
tf=tb(s,0.9,5.85,11.5,0.6)
line(tf,[("→ 외부 솔루션 구독이 아니라 ",13.5,MUTED,False),("우리 업무에 정확히 맞춘 자체 시스템.",13.5,INK,True),(" 필요한 기능을 바로 추가합니다.",13.5,MUTED,False)],first=True,after=0)
footer(s,"기술·보안",23)

# =========================================================
# SLIDE 21 — 로드맵
# =========================================================
s = slide(CREAM)
kicker(s, "06 · 향후 로드맵", y=0.55)
title(s, "지금도 돌아가고, 더 단단해집니다.", y=1.0, size=29)
road = [
    ("알림 강화",BLUE,"현재 앱 내 알림 → 카카오 알림톡·문자·웹푸시로 도달률 향상."),
    ("자동 검증",GREEN,"인스타 게시물 태그·공동작업자 자동 검증(Meta 연동). 현재는 자가 체크."),
    ("권한 분리",PURPLE,"팀원 개별 로그인·권한 (현재 공용 접속)."),
    ("데이터 고도화",PINK,"카테고리 세분화(뷰티→색조/스킨케어), 브랜드별 중복 전달 방지."),
    ("매출 관리",ORANGE,"캠페인 매출·순이익·마진과 브랜드 예산 연동."),
    ("고객 경험",RGBColor(0x55,0x50,0x5e),"FAQ·1:1 문의 고도화, 업로드 마감 자동 제어, 정산 주민번호 암호화."),
]
cw=3.73; ch=1.85; gx=0.18; gy=0.2; x0=0.9; y0=2.2
for i,(tg,col,desc) in enumerate(road):
    r=i//3; c=i%3
    x=x0+c*(cw+gx); y=y0+r*(ch+gy)
    card(s,x,y,cw,ch)
    # tag pill
    pw=0.34+len(tg)*0.12
    p=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x+0.2),Inches(y+0.2),Inches(pw),Inches(0.36))
    p.fill.solid(); p.fill.fore_color.rgb=col; p.line.fill.background(); p.shadow.inherit=False
    try: p.adjustments[0]=0.5
    except: pass
    ptf=p.text_frame; ptf.margin_top=0;ptf.margin_bottom=0; pp=ptf.paragraphs[0]; pp.alignment=PP_ALIGN.CENTER
    rr=pp.add_run(); rr.text=tg; rr.font.size=Pt(11.5); rr.font.bold=True; rr.font.color.rgb=WHITE; rr.font.name=F; _setfont(rr,F)
    tf=tb(s,x+0.2,y+0.72,cw-0.4,ch-0.85)
    line(tf,(desc,12.5,MUTED,False),first=True,lh=1.35,after=0)
footer(s,"로드맵",24)

# =========================================================
# SLIDE 22 — 마무리
# =========================================================
s = slide(DARK)
kicker(s, "SIRIAI", y=0.95, dark=True)
tf = tb(s, 1.0, 1.7, 11.3, 1.8, anchor=MSO_ANCHOR.MIDDLE)
line(tf, ("수작업으로는 못 키우던 규모를,", 38, DTEXT, True), first=True, align=PP_ALIGN.CENTER, after=4, lh=1.1)
line(tf, ("시스템으로 키웁니다.", 38, PEACH, True), align=PP_ALIGN.CENTER, after=0, lh=1.1)
eff = [("↓","매니저 반복 공수"),("↑","클라이언트 감도·만족"),("∞","데이터 자산화·확장성")]
cw=3.2; gap=0.3; total=cw*3+gap*2; x0=(SW-total)/2; cy=3.95
for i,(big,lab) in enumerate(eff):
    x=x0+i*(cw+gap)
    card(s,x,cy,cw,1.5,fill=DCARD,ln=DLINE)
    tf=tb(s,x,cy+0.2,cw,1.1,anchor=MSO_ANCHOR.MIDDLE)
    line(tf,(big,30,PEACH,True),first=True,align=PP_ALIGN.CENTER,after=2)
    line(tf,(lab,12,DMUTE,False),align=PP_ALIGN.CENTER,after=0)
tf=tb(s,1.0,5.85,11.3,0.5)
line(tf,("감사합니다.   ·   질문을 받겠습니다.",15,DTEXT,False),first=True,align=PP_ALIGN.CENTER,after=0)
footer(s,"마무리",25,dark=True)

# ---- save ----
out = r"C:\Users\User\Desktop\바이브코딩용 폴더\SIRIAI\SIRIAI_발표자료.pptx"
prs.save(out)
print("SAVED", out, "slides:", len(prs.slides._sldIdLst))
